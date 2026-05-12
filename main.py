import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

# --- Styling Constants ---
BG_DARK = "#1e1e1e"
BG_LIGHT = "#2d2d2d"
FG_TEXT = "#e0e0e0"
ACCENT_BLUE = "#3498db"
ACCENT_GREEN = "#27ae60"
ACCENT_GREY = "#555555"


class PPTXBuilderApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PPTX Definition Builder Pro")
        self.root.geometry("800x850")
        self.root.configure(bg=BG_DARK)

        self.txt_path = ""
        self.pptx_path = ""
        self.raw_topics = []
        self.current_index = 0

        self.setup_ui()

    def setup_ui(self):
        # --- File Selection Section ---
        file_frame = tk.LabelFrame(self.root, text=" 1. File Setup ", bg=BG_DARK, fg=ACCENT_BLUE, padx=10, pady=10)
        file_frame.pack(pady=10, padx=20, fill="x")

        self.btn_input = tk.Button(file_frame, text="Select Input (.txt)", command=self.select_input,
                                   bg=BG_LIGHT, fg=FG_TEXT, relief="flat", width=20)
        self.btn_input.grid(row=0, column=0, pady=5)
        self.lbl_input = tk.Label(file_frame, text="No file selected", bg=BG_DARK, fg="grey", anchor="w")
        self.lbl_input.grid(row=0, column=1, padx=10, sticky="w")

        self.btn_output = tk.Button(file_frame, text="Select Output (.pptx)", command=self.select_output,
                                    bg=BG_LIGHT, fg=FG_TEXT, relief="flat", width=20)
        self.btn_output.grid(row=1, column=0, pady=5)
        self.lbl_output = tk.Label(file_frame, text="No location selected", bg=BG_DARK, fg="grey", anchor="w")
        self.lbl_output.grid(row=1, column=1, padx=10, sticky="w")

        self.btn_start = tk.Button(file_frame, text="START SESSION", command=self.start_session,
                                   bg=ACCENT_BLUE, fg="white", state="disabled", font=("Arial", 10, "bold"),
                                   relief="flat")
        self.btn_start.grid(row=2, column=0, columnspan=2, pady=10, sticky="we")

        # --- Progress Section ---
        progress_frame = tk.Frame(self.root, bg=BG_DARK)
        progress_frame.pack(fill="x", padx=25)

        self.prog_bar = ttk.Progressbar(progress_frame, orient="horizontal", length=100, mode="determinate")
        self.prog_bar.pack(side="left", fill="x", expand=True, pady=5)

        self.lbl_percent = tk.Label(progress_frame, text="0%", bg=BG_DARK, fg=FG_TEXT, font=("Arial", 10, "bold"),
                                    width=6)
        self.lbl_percent.pack(side="right")

        # --- Main Input Section ---
        self.input_frame = tk.LabelFrame(self.root, text=" 2. Definitions ", bg=BG_DARK, fg=ACCENT_BLUE, padx=10,
                                         pady=10)
        self.input_frame.pack(pady=10, padx=20, fill="both", expand=True)

        self.topic_info_label = tk.Label(self.input_frame, text="Topic 0 of 0", bg=BG_DARK, fg=ACCENT_BLUE)
        self.topic_info_label.pack()

        # --- LIVE PREVIEW AREA ---
        self.preview_frame = tk.Frame(self.input_frame, bg="white", height=200, relief="sunken", bd=2)
        self.preview_frame.pack(fill="x", padx=10, pady=10)
        self.preview_frame.pack_propagate(False)  # Keep fixed height

        self.pre_cat = tk.Label(self.preview_frame, text="CATEGORY", font=("Arial", 10, "bold"), bg="white", fg="grey",
                                anchor="w")
        self.pre_cat.pack(fill="x", padx=20, pady=(10, 0))

        self.pre_topic = tk.Label(self.preview_frame, text="Slide Topic Title", font=("Arial", 18, "bold"), bg="white",
                                  fg="black", anchor="w")
        self.pre_topic.pack(fill="x", padx=20)

        self.pre_def = tk.Label(self.preview_frame, text="Definition preview will appear here...",
                                font=("Arial", 12), bg="white", fg="#333", wraplength=500, justify="left", anchor="nw")
        self.pre_def.pack(fill="both", expand=True, padx=20, pady=5)

        # --- Text Input ---
        self.text_box = tk.Text(self.input_frame, height=6, bg=BG_LIGHT, fg=FG_TEXT,
                                insertbackground="white", font=("Arial", 12), wrap="word", state="disabled", bd=0)
        self.text_box.pack(pady=10, padx=10, fill="x")
        self.text_box.bind("<KeyRelease>", self.update_preview)
        self.text_box.bind("<Return>", self.handle_enter)

        # --- Navigation ---
        nav_frame = tk.Frame(self.input_frame, bg=BG_DARK)
        nav_frame.pack(pady=10)

        self.back_btn = tk.Button(nav_frame, text="Back", command=self.go_back, state="disabled",
                                  bg=ACCENT_GREY, fg=FG_TEXT, width=10, relief="flat")
        self.back_btn.grid(row=0, column=0, padx=5)

        self.submit_btn = tk.Button(nav_frame, text="Submit & Save", command=self.save_current_step,
                                    bg=ACCENT_GREEN, fg="white", state="disabled", width=20, font=("Arial", 10, "bold"),
                                    relief="flat")
        self.submit_btn.grid(row=0, column=1, padx=5)

        self.skip_btn = tk.Button(nav_frame, text="Skip", command=self.skip_step, state="disabled",
                                  bg=ACCENT_GREY, fg=FG_TEXT, width=10, relief="flat")
        self.skip_btn.grid(row=0, column=2, padx=5)

    def update_preview(self, event=None):
        content = self.text_box.get("1.0", tk.END).strip()
        self.pre_def.config(text=content if content else "...")

    def select_input(self):
        path = filedialog.askopenfilename(filetypes=[("Text files", "*.txt")])
        if path:
            self.txt_path = path
            self.lbl_input.config(text=os.path.basename(path), fg=FG_TEXT)
            self.check_ready()

    def select_output(self):
        path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")])
        if path:
            self.pptx_path = path
            self.lbl_output.config(text=os.path.basename(path), fg=FG_TEXT)
            self.check_ready()

    def check_ready(self):
        if self.txt_path and self.pptx_path:
            self.btn_start.config(state="normal")

    def start_session(self):
        self.raw_topics = []
        try:
            with open(self.txt_path, "r", encoding="utf-8") as f:
                curr_cat = "General"
                for line in f:
                    line = line.strip()
                    if not line: continue
                    if line.startswith("●"):
                        self.raw_topics.append({"cat": curr_cat, "topic": line.lstrip("● ").strip()})
                    else:
                        curr_cat = line
        except Exception as e:
            messagebox.showerror("Error", f"Read Error: {e}")
            return

        self.current_index = 0
        if os.path.exists(self.pptx_path):
            try:
                prs = Presentation(self.pptx_path)
                existing_names = [s.name for s in prs.slides]
                while self.current_index < len(self.raw_topics):
                    if self.raw_topics[self.current_index]['topic'] in existing_names:
                        self.current_index += 1
                    else:
                        break
            except:
                pass

        self.text_box.config(state="normal")
        self.submit_btn.config(state="normal")
        self.skip_btn.config(state="normal")
        self.update_display()
        self.text_box.focus_force()

    def update_display(self):
        total = len(self.raw_topics)
        if total == 0: return

        self.back_btn.config(state="normal" if self.current_index > 0 else "disabled")

        # Update Progress
        percent = int((self.current_index / total) * 100)
        self.prog_bar['value'] = percent
        self.lbl_percent.config(text=f"{percent}%")

        if self.current_index < total:
            item = self.raw_topics[self.current_index]
            self.topic_info_label.config(text=f"Topic {self.current_index + 1} of {total}")

            # Update Preview static fields
            self.pre_cat.config(text=item['cat'].upper())
            self.pre_topic.config(text=item['topic'])

            # Load existing text if any
            self.text_box.delete("1.0", tk.END)
            existing = self.get_existing_text(item['topic'])
            if existing:
                self.text_box.insert("1.0", existing)

            self.update_preview()
        else:
            messagebox.showinfo("Finished", "All topics completed!")
            self.root.destroy()

    def get_existing_text(self, topic_name):
        if not os.path.exists(self.pptx_path): return None
        try:
            prs = Presentation(self.pptx_path)
            for slide in prs.slides:
                if slide.name == topic_name and len(slide.shapes) >= 3:
                    return slide.shapes[2].text_frame.text
            return None
        except:
            return None

    def handle_enter(self, event):
        if not (event.state & 0x1):  # If Shift not held
            self.save_current_step()
            return "break"

    def save_current_step(self):
        definition = self.text_box.get("1.0", tk.END).strip()
        if not definition:
            messagebox.showwarning("Empty", "Please provide a definition.")
            return

        item = self.raw_topics[self.current_index]
        self.add_or_update_slide(item['cat'], item['topic'], definition)
        self.current_index += 1
        self.update_display()

    def go_back(self):
        if self.current_index > 0:
            self.current_index -= 1
            self.update_display()

    def skip_step(self):
        self.current_index += 1
        self.update_display()

    def add_or_update_slide(self, category, topic, definition):
        try:
            prs = Presentation(self.pptx_path) if os.path.exists(self.pptx_path) else Presentation()
        except:
            prs = Presentation()

        target_slide = next((s for s in prs.slides if s.name == topic), None)

        if target_slide:
            if len(target_slide.shapes) >= 3:
                target_slide.shapes[0].text_frame.text = category.upper()
                target_slide.shapes[1].text_frame.text = topic
                target_slide.shapes[2].text_frame.text = definition
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.name = topic
            # Formatting logic remains the same as your original
            for i, (y, h, sz, bld, txt) in enumerate([
                (0.2, 0.5, 14, True, category.upper()),
                (0.7, 0.8, 28, True, topic),
                (1.6, 5.4, 20, False, definition)
            ]):
                shape = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(h))
                frame = shape.text_frame
                frame.word_wrap = True
                if i == 2: frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                p = frame.paragraphs[0]
                p.text = txt
                p.font.size = Pt(sz)
                p.font.bold = bld

        try:
            prs.save(self.pptx_path)
        except PermissionError:
            messagebox.showerror("File Error", "Please close the PPTX file first!")


if __name__ == "__main__":
    root = tk.Tk()
    # Simple style adjustments for the progress bar to match dark mode
    style = ttk.Style()
    style.theme_use('default')
    style.configure("TProgressbar", thickness=10, background=ACCENT_BLUE, troughcolor=BG_LIGHT, bordercolor=BG_DARK)

    app = PPTXBuilderApp(root)
    root.mainloop()